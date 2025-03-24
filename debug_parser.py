"""
直接测试TemplateParser类的脚本
"""
import os
import sys
import logging
from pptx import Presentation

# 配置日志
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# 添加backend目录到路径
sys.path.insert(0, os.path.abspath('backend'))

# 创建测试PPT
def create_test_ppt():
    logger.info("创建测试PPT文件...")
    prs = Presentation()
    
    # 添加标题页
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "测试标题"
    subtitle.text = "测试副标题"
    
    test_ppt = "test_debug.pptx"
    prs.save(test_ppt)
    logger.info(f"创建测试PPT文件: {test_ppt}")
    return os.path.abspath(test_ppt)

def test_parser(file_path):
    from core.template_parser import TemplateParser
    
    try:
        logger.info(f"尝试解析文件: {file_path}")
        parser = TemplateParser(file_path)
        template = parser.parse()
        
        if template:
            logger.info("解析成功!")
            logger.info(f"Template ID: {template.template_id}")
            logger.info(f"Name: {template.name}")
            logger.info(f"Slides: {len(template.slides)}")
            for i, slide in enumerate(template.slides):
                logger.info(f"Slide {i+1}: ID={slide.slide_id}, Type={slide.slide_type}, Elements={len(slide.elements)}")
            return True
        else:
            logger.error("解析失败，返回了None")
            return False
    except Exception as e:
        logger.exception(f"解析时出错: {e}")
        return False

def main():
    try:
        # 创建测试PPT
        test_ppt = create_test_ppt()
        
        # 测试解析
        success = test_parser(test_ppt)
        if success:
            logger.info("模板解析成功!")
        else:
            logger.error("模板解析失败!")
        
        # 清理测试文件
        # os.remove(test_ppt)
        # logger.info(f"已删除测试文件: {test_ppt}")
    except Exception as e:
        logger.exception(f"测试过程中出错: {e}")

if __name__ == "__main__":
    main() 