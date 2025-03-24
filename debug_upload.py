"""
调试模板上传功能的脚本
"""
import requests
import os
import sys
from pptx import Presentation
import logging

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# API基础URL
BASE_URL = "http://localhost:8000"

def create_test_ppt():
    """创建测试PPT文件"""
    logger.info("创建测试PPT文件...")
    
    # 创建一个简单的测试PPT
    prs = Presentation()
    
    # 添加标题页
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "测试标题"
    subtitle.text = "测试副标题"
    
    test_ppt = "test_template.pptx"
    prs.save(test_ppt)
    logger.info(f"创建测试PPT文件: {test_ppt}")
    return test_ppt

def upload_template(file_path):
    """上传模板文件"""
    logger.info(f"上传模板文件: {file_path}")
    
    try:
        with open(file_path, "rb") as f:
            files = {"file": (os.path.basename(file_path), f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
            response = requests.post(f"{BASE_URL}/templates/upload", files=files)
            
            logger.info(f"状态码: {response.status_code}")
            if response.status_code == 200:
                logger.info(f"响应: {response.json()}")
                return response.json()
            else:
                logger.error(f"错误响应: {response.text}")
                return None
    except Exception as e:
        logger.error(f"上传模板文件失败: {str(e)}")
        return None

def main():
    """主函数"""
    try:
        # 创建测试PPT
        test_ppt = create_test_ppt()
        
        # 上传模板
        result = upload_template(test_ppt)
        if result:
            logger.info("模板上传成功！")
        else:
            logger.error("模板上传失败！")
        
        # 清理测试文件
        # os.remove(test_ppt)
        # logger.info(f"已删除测试文件: {test_ppt}")
    except Exception as e:
        logger.error(f"测试过程中出错: {str(e)}")

if __name__ == "__main__":
    main() 