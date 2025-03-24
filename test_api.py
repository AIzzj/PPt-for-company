"""
API功能测试脚本
"""
import os
import sys
import logging
import requests
import json
from pptx import Presentation
from time import sleep

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# API基础URL
BASE_URL = "http://localhost:8000"

def test_root():
    """测试根路由"""
    logger.info("测试根路由...")
    try:
        response = requests.get(f"{BASE_URL}/")
        response.raise_for_status()
        logger.info(f"状态码: {response.status_code}")
        logger.info(f"响应: {response.json()}")
        return True
    except requests.exceptions.RequestException as e:
        logger.error(f"测试根路由失败: {str(e)}")
        return False

def test_template_upload():
    """测试模板上传"""
    logger.info("测试模板上传...")
    try:
        # 创建测试目录
        os.makedirs("test_files", exist_ok=True)
        
        # 创建一个简单的测试PPT
        prs = Presentation()
        
        # 添加标题页
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "测试标题"
        subtitle.text = "测试副标题"
        
        # 添加内容页
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = content_slide.shapes.title
        content = content_slide.placeholders[1]
        title.text = "内容页标题"
        content.text = "• 测试内容1\n• 测试内容2\n• 测试内容3"
        
        test_ppt = "test_files/test_template.pptx"
        prs.save(test_ppt)
        logger.info(f"创建测试PPT文件: {test_ppt}")
        
        # 上传模板
        with open(test_ppt, "rb") as f:
            files = {"file": ("test_template.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
            response = requests.post(f"{BASE_URL}/templates/upload", files=files)
            response.raise_for_status()
            
        logger.info(f"状态码: {response.status_code}")
        logger.info(f"响应: {response.json()}")
        return response.json().get("template_id")
    except requests.exceptions.RequestException as e:
        logger.error(f"上传模板失败: {str(e)}")
        if hasattr(e, 'response') and e.response is not None and hasattr(e.response, 'json'):
            try:
                logger.error(f"错误详情: {e.response.json()}")
            except:
                logger.error("无法解析错误响应内容")
        return None
    except Exception as e:
        logger.error(f"创建或上传模板时出错: {str(e)}")
        return None

def test_list_templates():
    """测试列出模板"""
    logger.info("测试列出模板...")
    try:
        response = requests.get(f"{BASE_URL}/templates")
        response.raise_for_status()
        logger.info(f"状态码: {response.status_code}")
        logger.info(f"响应: {response.json()}")
        return True
    except requests.exceptions.RequestException as e:
        logger.error(f"列出模板失败: {str(e)}")
        return False

def test_generate_content(template_id):
    """测试内容生成"""
    logger.info("测试内容生成...")
    try:
        data = {
            "template_id": template_id,
            "slide_id": "slide_1",
            "element_id": "title_1",
            "content_type": "title",
            "user_input": "生成一个创新的标题",
            "style_guide": {"tone": "专业", "length": "简短"}
        }
        response = requests.post(f"{BASE_URL}/content/generate", json=data)
        response.raise_for_status()
        logger.info(f"状态码: {response.status_code}")
        logger.info(f"响应: {response.json()}")
        return response.json()
    except requests.exceptions.RequestException as e:
        logger.error(f"生成内容失败: {str(e)}")
        if hasattr(e, 'response') and e.response is not None and hasattr(e.response, 'json'):
            try:
                logger.error(f"错误详情: {e.response.json()}")
            except:
                logger.error("无法解析错误响应内容")
        return None

def test_generate_ppt(template_id, content_response):
    """测试PPT生成"""
    logger.info("测试PPT生成...")
    try:
        data = {
            "template_id": template_id,
            "content_mappings": [{
                "slide_id": "slide_1",
                "element_id": "title_1",
                "content": content_response["content"],
                "content_type": content_response["content_type"]
            }]
        }
        response = requests.post(f"{BASE_URL}/ppt/generate", json=data)
        response.raise_for_status()
        logger.info(f"状态码: {response.status_code}")
        logger.info(f"响应: {response.json()}")
        return response.json().get("task_id")
    except requests.exceptions.RequestException as e:
        logger.error(f"生成PPT失败: {str(e)}")
        if hasattr(e, 'response') and e.response is not None and hasattr(e.response, 'json'):
            try:
                logger.error(f"错误详情: {e.response.json()}")
            except:
                logger.error("无法解析错误响应内容")
        return None

def test_get_status(task_id):
    """测试获取任务状态"""
    logger.info("测试获取任务状态...")
    try:
        max_retries = 5
        for i in range(max_retries):
            response = requests.get(f"{BASE_URL}/ppt/status/{task_id}")
            response.raise_for_status()
            status_data = response.json()
            logger.info(f"状态码: {response.status_code}")
            logger.info(f"响应: {status_data}")
            
            if status_data["status"] == "completed":
                return True
            elif i < max_retries - 1:
                logger.info("任务仍在进行中，等待5秒后重试...")
                sleep(5)
        
        logger.warning("达到最大重试次数，任务可能仍在进行中")
        return False
    except requests.exceptions.RequestException as e:
        logger.error(f"获取任务状态失败: {str(e)}")
        if hasattr(e, 'response') and e.response is not None and hasattr(e.response, 'json'):
            try:
                logger.error(f"错误详情: {e.response.json()}")
            except:
                logger.error("无法解析错误响应内容")
        return False

def main():
    """运行所有测试"""
    success = True
    
    # 测试根路由
    if not test_root():
        logger.error("根路由测试失败")
        success = False
        return
    
    # 测试模板上传
    template_id = test_template_upload()
    if not template_id:
        logger.error("模板上传测试失败")
        success = False
        return
    
    # 测试列出模板
    if not test_list_templates():
        logger.error("列出模板测试失败")
        success = False
        return
    
    # 测试内容生成
    content_response = test_generate_content(template_id)
    if not content_response:
        logger.error("内容生成测试失败")
        success = False
        return
    
    # 测试PPT生成
    task_id = test_generate_ppt(template_id, content_response)
    if not task_id:
        logger.error("PPT生成测试失败")
        success = False
        return
    
    # 测试任务状态
    if not test_get_status(task_id):
        logger.error("获取任务状态测试失败")
        success = False
        return
    
    if success:
        logger.info("所有测试完成")
    else:
        logger.error("测试过程中出现错误")
        sys.exit(1)

if __name__ == "__main__":
    main() 